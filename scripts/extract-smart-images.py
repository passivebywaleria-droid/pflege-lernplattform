#!/usr/bin/env python3
"""
Extrahiert alle Bilder aus SMART PPTX-Dateien nach public/images/.
Jede PPTX → eigener Unterordner mit nummerierten PNGs.
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = Path(__file__).parent.parent / "public" / "images"
CATEGORIES = ["anatomy", "pathology", "equipment"]

def extract_images_from_pptx(pptx_path: Path, output_dir: Path):
    """Extract all images from a PPTX file into output_dir."""
    if pptx_path.stat().st_size < 50_000:  # Skip files < 50KB (likely broken)
        print(f"  SKIP (nur {pptx_path.stat().st_size // 1024}K — wahrscheinlich kaputt)")
        return 0

    prs = Presentation(str(pptx_path))
    output_dir.mkdir(parents=True, exist_ok=True)

    img_count = 0
    seen_blobs = set()

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # Direct image
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                blob_hash = hash(blob)
                if blob_hash in seen_blobs:
                    continue
                seen_blobs.add(blob_hash)

                ext = shape.image.content_type.split("/")[-1]
                if ext == "x-emf" or ext == "x-wmf":
                    ext = "emf" if "emf" in ext else "wmf"
                if ext not in ("png", "jpeg", "jpg", "gif", "svg+xml", "emf", "wmf", "tiff"):
                    ext = "png"
                if ext == "jpeg":
                    ext = "jpg"
                if ext == "svg+xml":
                    ext = "svg"

                img_count += 1
                filename = f"{img_count:03d}.{ext}"
                filepath = output_dir / filename
                with open(filepath, "wb") as f:
                    f.write(blob)

            # Group shapes may contain images
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    if subshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        blob = subshape.image.blob
                        blob_hash = hash(blob)
                        if blob_hash in seen_blobs:
                            continue
                        seen_blobs.add(blob_hash)

                        ext = subshape.image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        if ext == "svg+xml":
                            ext = "svg"
                        if ext not in ("png", "jpg", "gif", "svg", "emf", "wmf", "tiff"):
                            ext = "png"

                        img_count += 1
                        filename = f"{img_count:03d}.{ext}"
                        filepath = output_dir / filename
                        with open(filepath, "wb") as f:
                            f.write(blob)

    return img_count


def main():
    total = 0
    catalog_lines = []

    for category in CATEGORIES:
        cat_dir = BASE / category
        if not cat_dir.exists():
            continue

        for pptx_file in sorted(cat_dir.glob("*.pptx")):
            # Create subfolder from filename: SMART-Heart-physiology.pptx → heart-physiology
            folder_name = pptx_file.stem.replace("SMART-", "").lower()
            output_dir = cat_dir / folder_name

            print(f"\n{category}/{pptx_file.name}:")
            count = extract_images_from_pptx(pptx_file, output_dir)
            print(f"  → {count} Bilder extrahiert nach {category}/{folder_name}/")
            total += count

            if count > 0:
                catalog_lines.append(f"| {category}/{folder_name}/ | {count} | {', '.join(f.name for f in sorted(output_dir.iterdir())[:5])}{'...' if count > 5 else ''} |")

    print(f"\n{'='*60}")
    print(f"GESAMT: {total} Bilder extrahiert")
    print(f"{'='*60}")

    # Write catalog summary
    catalog_path = BASE.parent.parent / "content" / "_bild-katalog-raw.md"
    with open(catalog_path, "w") as f:
        f.write(f"# SMART Bild-Katalog (auto-generiert)\n\n")
        f.write(f"Gesamt: {total} Bilder\n\n")
        f.write("| Ordner | Anzahl | Beispiele |\n")
        f.write("|--------|--------|-----------|\n")
        for line in catalog_lines:
            f.write(line + "\n")

    return total


if __name__ == "__main__":
    count = main()
    sys.exit(0 if count > 0 else 1)
